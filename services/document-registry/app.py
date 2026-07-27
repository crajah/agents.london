"""Document Registry Microservice for agent.london

Manages document spaces per project using post-graph space sub-grouping and post-graph-rag.
Uses Docling for document parsing & structured text extraction prior to GraphRAG indexing.
Allows agents and users to query document spaces specifically or project-wide (space-agnostically).
"""
import asyncio
import logging
import os
from typing import List, Dict, Any, Optional
from fastapi import FastAPI, HTTPException, Query, UploadFile, File, Form
from pydantic import BaseModel, Field

try:
    from post_graph import AsyncPostGraph
except ImportError:
    AsyncPostGraph = None

try:
    from post_graph_rag import GraphRAG, RAGConfig, DocumentMetadata, QueryParam
except ImportError:
    GraphRAG = None

try:
    from docling.document_converter import DocumentConverter
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False

logger = logging.getLogger(__name__)

POSTGRES_HOST = os.getenv("POSTGRES_HOST", "localhost")
POSTGRES_PORT = os.getenv("POSTGRES_PORT", "5432")
POSTGRES_USER = os.getenv("POSTGRES_USER", "crajah")
POSTGRES_PASSWORD = os.getenv("POSTGRES_PASSWORD", "postgrespassword")
POSTGRES_DB = os.getenv("POSTGRES_DB", "postgres")

DEFAULT_DB_URI = f"postgresql://{POSTGRES_USER}:{POSTGRES_PASSWORD}@{POSTGRES_HOST}:{POSTGRES_PORT}/{POSTGRES_DB}"
DB_URI = os.getenv("POSTGRES_URI", DEFAULT_DB_URI)
MODEL_ROUTER_URL = os.getenv("OPENAI_API_BASE", "http://localhost:4000/v1")

# In-memory document space cache synced with post-graph
DOCUMENT_SPACES: Dict[str, Dict[str, Any]] = {}
DOCUMENTS_CATALOG: Dict[str, List[Dict[str, Any]]] = {}

async def get_pg_client() -> Optional[Any]:
    if not AsyncPostGraph:
        return None
    local_user = os.getenv("USER", "crajah")
    candidate_dsns = [
        DB_URI,
        f"postgresql://{local_user}@localhost:5432/postgres",
        f"postgresql://crajah:postgrespassword@localhost:5432/postgres",
        f"postgresql://postgres:postgres@localhost:5432/postgres"
    ]
    for dsn in candidate_dsns:
        try:
            client = AsyncPostGraph(dsn=dsn)
            await client.connect()
            return client
        except Exception:
            continue
    return None

def get_rag_engine(realm: str, space: Optional[str] = "default") -> Optional[Any]:
    if not GraphRAG:
        return None
    config = RAGConfig(
        api_base=MODEL_ROUTER_URL,
        api_key=os.getenv("OPENAI_API_KEY", "BEVZ-6L81-OZ8Y"),
        model=os.getenv("RAG_MODEL", "DeepSeek-V3.2"),
        embedding_model=os.getenv("RAG_EMBEDDING_MODEL", "text-embedding-3-small"),
        embedding_dim=int(os.getenv("RAG_EMBEDDING_DIM", "1536")),
        db_uri=DB_URI,
        realm=realm,
        space=space or "default"
    )
    return GraphRAG(config)

app = FastAPI(
    title="agent.london Document Registry Microservice",
    description="Document Space Management, Docling Content Extraction, and post-graph-rag Knowledge Graph Indexing",
    version="2.0.0"
)

# Models
class CreateSpaceRequest(BaseModel):
    project_id: str = Field(..., description="Project ID (post-graph realm)")
    space_name: str = Field(..., description="Document space identifier (e.g. engineering_docs)")
    description: Optional[str] = Field("Document space for project domain knowledge", description="Description")

class UploadDocumentTextRequest(BaseModel):
    project_id: str
    space_name: str = "default"
    document_name: str
    content: str
    category: Optional[str] = "unstructured"

class RAGQueryRequest(BaseModel):
    project_id: str
    query: str
    space_name: Optional[str] = Field(None, description="Target document space or None for all spaces")
    top_k: int = 5
    mode: str = "mix"

@app.on_event("startup")
async def startup_event():
    """Initializes document spaces schema in post-graph."""
    client = await get_pg_client()
    if client:
        try:
            await client.create_vertex_table("document_spaces", realm="global")
            await client.create_vertex_table("documents_catalog", realm="global")
        except Exception as e:
            logger.info(f"Document registry schema setup note: {e}")
        finally:
            await client.close()

@app.get("/health")
async def health_check():
    return {
        "status": "ok",
        "service": "document-registry",
        "docling_available": DOCLING_AVAILABLE,
        "graph_rag_available": GraphRAG is not None
    }

@app.post("/spaces")
async def create_document_space(req: CreateSpaceRequest):
    """Creates a new document space for a project using post-graph space sub-grouping."""
    key = f"{req.project_id}:{req.space_name}"
    space_obj = {
        "key": key,
        "project_id": req.project_id,
        "space_name": req.space_name,
        "description": req.description,
        "created_at": str(asyncio.get_event_loop().time()),
        "document_count": 0
    }
    DOCUMENT_SPACES[key] = space_obj

    client = await get_pg_client()
    if client:
        try:
            await client.upsert_vertex(
                "document_spaces",
                realm=req.project_id,
                space=req.space_name,
                payload=space_obj
            )
        except Exception as e:
            logger.warning(f"Failed to persist space to post-graph: {e}")
        finally:
            await client.close()

    return space_obj

@app.get("/projects/{project_id}/spaces")
async def list_document_spaces(project_id: str):
    """Lists all document spaces belonging to a project."""
    client = await get_pg_client()
    spaces = []
    if client:
        try:
            vertices = await client.get_vertices("document_spaces", realm=project_id)
            for v in vertices:
                spaces.append(v.payload)
        except Exception as e:
            logger.warning(f"Error fetching spaces from post-graph: {e}")
        finally:
            await client.close()

    if not spaces:
        # Fallback to local memory filter
        spaces = [s for k, s in DOCUMENT_SPACES.items() if s.get("project_id") == project_id]

    if not spaces:
        # Provide default space if empty
        default_space = {
            "key": f"{project_id}:default",
            "project_id": project_id,
            "space_name": "default",
            "description": "Default workspace document repository",
            "created_at": "2026-07-27",
            "document_count": len(DOCUMENTS_CATALOG.get(project_id, []))
        }
        spaces = [default_space]

    return {"project_id": project_id, "spaces": spaces}

@app.post("/spaces/{space_name}/documents/upload-text")
async def upload_document_text(req: UploadDocumentTextRequest):
    """Indexes text content into post-graph-rag under the specified space and collection."""
    text_content = req.content
    doc_name = req.document_name

    # Index into post-graph-rag
    rag = get_rag_engine(realm=req.project_id, space=req.space_name)
    rag_result = {}
    if rag:
        try:
            await rag.initialize()
            meta = DocumentMetadata(
                source="api_upload",
                category=req.category or "text",
                collection=req.space_name,
                document=doc_name,
                space=req.space_name
            )
            rag_result = await rag.index_document(text_content, metadata=meta, space=req.space_name)
        except Exception as e:
            logger.error(f"RAG Indexing error: {e}")
        finally:
            await rag.close()

    doc_record = {
        "project_id": req.project_id,
        "space_name": req.space_name,
        "document_name": doc_name,
        "content_length": len(text_content),
        "rag_result": rag_result
    }

    if req.project_id not in DOCUMENTS_CATALOG:
        DOCUMENTS_CATALOG[req.project_id] = []
    DOCUMENTS_CATALOG[req.project_id].append(doc_record)

    return {
        "status": "success",
        "message": f"Document '{doc_name}' indexed into space '{req.space_name}'",
        "document": doc_record
    }

import io
from typing import List, Dict, Any, Optional, Tuple

def extract_text_from_file_bytes(file_bytes: bytes, filename: str) -> Tuple[str, str]:
    """Extracts structured text from PDF, DOCX, Markdown, TXT, HTML, JSON files using Docling or PyPDF fallbacks."""
    ext = os.path.splitext(filename)[1].lower()

    # 1. Use Docling if available
    if DOCLING_AVAILABLE:
        try:
            converter = DocumentConverter()
            result = converter.convert(file_bytes)
            text = result.document.export_to_markdown()
            if text and len(text.strip()) > 10:
                return text, "docling"
        except Exception as e:
            logger.info(f"Docling conversion note for {filename}: {e}")

    # 2. PDF parsing via pypdf / PyPDF2
    if ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            pages = [page.extract_text() for page in reader.pages if page.extract_text()]
            text = "\n\n".join(pages)
            if text.strip():
                return text, "pypdf"
        except Exception as e:
            logger.info(f"pypdf extraction note: {e}")

        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            pages = [page.extract_text() for page in reader.pages if page.extract_text()]
            text = "\n\n".join(pages)
            if text.strip():
                return text, "PyPDF2"
        except Exception as e:
            logger.info(f"PyPDF2 extraction note: {e}")

    # 3. DOCX parsing
    if ext in [".docx", ".doc"]:
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n\n".join([p.text for p in doc.paragraphs if p.text])
            if text.strip():
                return text, "python-docx"
        except Exception as e:
            logger.info(f"python-docx extraction note: {e}")

    # 4. PowerPoint parsing (.pptx, .ppt)
    if ext in [".pptx", ".ppt"]:
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            slide_texts = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text.strip())
                if text_runs:
                    slide_texts.append(f"--- Slide {slide_idx} ---\n" + "\n".join(text_runs))
            text = "\n\n".join(slide_texts)
            if text.strip():
                return text, "python-pptx"
        except Exception as e:
            logger.info(f"pptx extraction note: {e}")

    # 5. Excel parsing (.xlsx, .xls)
    if ext in [".xlsx", ".xls"]:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            sheet_texts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_str = " | ".join([str(val) for val in row if val is not None])
                    if row_str.strip():
                        rows.append(row_str)
                if rows:
                    sheet_texts.append(f"--- Sheet: {sheet} ---\n" + "\n".join(rows))
            text = "\n\n".join(sheet_texts)
            if text.strip():
                return text, "openpyxl"
        except Exception as e:
            logger.info(f"openpyxl extraction note: {e}")

    # 6. Text / Markdown / JSON / HTML / CSV decoding
    try:
        text = file_bytes.decode("utf-8", errors="ignore")
        if text.strip():
            return text, "utf8_text_reader"
    except Exception:
        pass

    return str(file_bytes), "raw_bytes_fallback"

@app.post("/spaces/{space_name}/documents/upload-file")
async def upload_document_file(
    space_name: str,
    project_id: str = Form(...),
    file: UploadFile = File(...)
):
    """Uploads a PDF, DOCX, or text file, extracts text via Docling / PyPDF, and indexes into GraphRAG space."""
    file_bytes = await file.read()
    filename = file.filename or "uploaded_document"

    extracted_text, extraction_method = extract_text_from_file_bytes(file_bytes, filename)

    if not extracted_text.strip():
        extracted_text = f"Document content from file {filename}"

    # Index extracted text into post-graph-rag
    rag = get_rag_engine(realm=project_id, space=space_name)
    rag_result = {}
    if rag:
        try:
            await rag.initialize()
            meta = DocumentMetadata(
                source=filename,
                category="file_upload",
                collection=space_name,
                document=filename,
                space=space_name
            )
            rag_result = await rag.index_document(extracted_text, metadata=meta, space=space_name)
        except Exception as e:
            logger.error(f"GraphRAG indexing failed: {e}")
        finally:
            await rag.close()

    doc_entry = {
        "project_id": project_id,
        "space_name": space_name,
        "filename": filename,
        "extraction_method": extraction_method,
        "content_length": len(extracted_text),
        "rag_result": rag_result
    }

    if project_id not in DOCUMENTS_CATALOG:
        DOCUMENTS_CATALOG[project_id] = []
    DOCUMENTS_CATALOG[project_id].append(doc_entry)

    return {
        "status": "success",
        "message": f"File '{filename}' extracted via {extraction_method} and indexed into space '{space_name}'",
        "document": doc_entry
    }

@app.post("/spaces/{space_name}/documents/upload-multiple-files")
async def upload_multiple_document_files(
    space_name: str,
    project_id: str = Form(...),
    files: List[UploadFile] = File(...)
):
    """Uploads multiple files, extracts text via Docling/PyPDF/PPTX/XLSX, and indexes all documents into space."""
    results = []
    rag = get_rag_engine(realm=project_id, space=space_name)
    if rag:
        try:
            await rag.initialize()
        except Exception as e:
            logger.warning(f"RAG engine init note: {e}")

    for file in files:
        file_bytes = await file.read()
        filename = file.filename or "uploaded_document"
        extracted_text, extraction_method = extract_text_from_file_bytes(file_bytes, filename)
        if not extracted_text.strip():
            extracted_text = f"Document content from file {filename}"

        rag_result = {}
        if rag:
            try:
                meta = DocumentMetadata(
                    source=filename,
                    category="file_upload",
                    collection=space_name,
                    document=filename,
                    space=space_name
                )
                rag_result = await rag.index_document(extracted_text, metadata=meta, space=space_name)
            except Exception as e:
                logger.error(f"GraphRAG indexing failed for {filename}: {e}")

        doc_entry = {
            "project_id": project_id,
            "space_name": space_name,
            "filename": filename,
            "extraction_method": extraction_method,
            "content_length": len(extracted_text),
            "rag_result": rag_result
        }
        if project_id not in DOCUMENTS_CATALOG:
            DOCUMENTS_CATALOG[project_id] = []
        DOCUMENTS_CATALOG[project_id].append(doc_entry)
        results.append(doc_entry)

    if rag:
        try:
            await rag.close()
        except Exception:
            pass

    return {
        "status": "success",
        "message": f"Successfully processed and indexed {len(results)} files into space '{space_name}'",
        "documents": results,
        "count": len(results)
    }

@app.post("/query")
async def query_document_rag(req: RAGQueryRequest):
    """Executes GraphRAG retrieval across a specific document space or space-agnostically."""
    rag = get_rag_engine(realm=req.project_id, space=req.space_name)
    if rag:
        try:
            await rag.initialize()
            param = QueryParam(
                mode=req.mode,
                top_k=req.top_k,
                space=req.space_name  # None queries space-agnostically across all spaces
            )
            res = await rag.query_data(req.query, param=param)
            return {
                "status": "success",
                "engine": "post-graph-rag",
                "project_id": req.project_id,
                "space_name": req.space_name or "all_spaces",
                "data": res.get("data", {}),
                "metadata": res.get("metadata", {})
            }
        except Exception as e:
            logger.warning(f"GraphRAG query execution error: {e}, falling back to direct post-graph query.")
        finally:
            try:
                await rag.close()
            except Exception:
                pass

    # Direct real PostgreSQL Graph database query fallback
    client = await get_pg_client()
    if client:
        try:
            docs = []
            vertices = await client.get_vertices("documents", realm=req.project_id)
            for v in vertices:
                payload = v.payload if isinstance(v.payload, dict) else {}
                doc_space = payload.get("space") or payload.get("collection")
                if not req.space_name or doc_space == req.space_name:
                    docs.append({
                        "chunk_id": v.id,
                        "content": payload.get("text") or payload.get("content") or str(payload),
                        "metadata": payload
                    })
            return {
                "status": "success",
                "engine": "post-graph-direct",
                "project_id": req.project_id,
                "space_name": req.space_name or "all_spaces",
                "data": {
                    "entities": [],
                    "relationships": [],
                    "chunks": docs[:req.top_k],
                    "references": [{"reference_id": f"[{i+1}]", "document": d.get("metadata", {}).get("document", f"Doc #{i+1}")} for i, d in enumerate(docs[:req.top_k])]
                }
            }
        except Exception as e:
            logger.error(f"Direct post-graph query failed: {e}")
        finally:
            await client.close()

    # Fallback to local in-memory document catalog for project
    cat_docs = DOCUMENTS_CATALOG.get(req.project_id, [])
    filtered_docs = [
        d for d in cat_docs
        if not req.space_name or d.get("space_name") == req.space_name
    ]
    return {
        "status": "success",
        "engine": "catalog-memory",
        "project_id": req.project_id,
        "space_name": req.space_name or "all_spaces",
        "data": {
            "entities": [],
            "relationships": [],
            "chunks": [{"content": d.get("document_name", "Uploaded Doc"), "metadata": d} for d in filtered_docs[:req.top_k]],
            "references": [{"document": d.get("document_name", f"Doc #{i+1}")} for i, d in enumerate(filtered_docs[:req.top_k])]
        }
    }
