"""Extracting text from an uploaded file (spec §5).

The ladder is tried in order and the first real success wins. Docling is first
because it produces *structured* Markdown — headings, lists, tables — and
structure survives chunking. A table flattened by a naive PDF text extractor
becomes a run of numbers with no column headers, which embeds as noise and
retrieves as nonsense.

**Memory discipline (2026-09-05).** Extraction is streamed: the upload is
spooled to a temp file by the route and every parser reads from the *path*,
so the raw bytes are never held whole in this process. Two OOM kills taught
the old shape's lesson — a full `bytes` copy per file, plus a fresh docling
``DocumentConverter`` loading its layout/OCR models on every call, meant one
real batch upload could exceed any limit we set. The converter is now a
module singleton (models load once and stay), and only the extracted *text*
— the product — lives in memory.

**The rule this module exists to enforce is Rule 5.3: when nothing succeeds,
extraction fails.** It does not substitute a placeholder. That rule has two
violations in its own history, which is why it is stated this emphatically:

  - the first returned `str(file_bytes)` — the Python repr of the raw bytes —
    as the document's text, which was embedded and indexed into the knowledge
    graph as if it were prose;
  - the second substituted the sentence "Document content from file {name}".

Both produce a corpus entry that retrieves, cites and reads as a real document,
and neither leaves any signal that extraction failed.
"""
from __future__ import annotations

import codecs
import logging
import os
import tempfile
from typing import Optional

from doc_model import DocumentError, ExtractionResult

logger = logging.getLogger(__name__)

# A parser that returns whitespace or a handful of characters has not succeeded;
# the ladder continues (Rule 5.1).
MIN_CHARACTERS = 10

_docling_converter = None


def _converter():
    """One DocumentConverter for the process lifetime. Its layout and OCR
    models are the largest allocation in this service; loading them per
    call turned every upload into a memory spike."""
    global _docling_converter
    if _docling_converter is None:
        from docling.document_converter import DocumentConverter
        _docling_converter = DocumentConverter()
    return _docling_converter


def _ok(text: Optional[str]) -> bool:
    return bool(text and len(text.strip()) >= MIN_CHARACTERS)


def _result(text: str, method: str) -> ExtractionResult:
    return ExtractionResult(method=method, text=text, characters=len(text))


def extract_path(path: str, filename: str) -> ExtractionResult:
    """Extract text from a spooled file, or raise DocumentError (Rule 5.3)."""
    ext = os.path.splitext(filename)[1].lower()

    for attempt in (_docling, _pdf, _docx, _pptx, _xlsx, _plain_text):
        try:
            found = attempt(path, ext)
        except Exception as e:
            # A parser raising is a normal rung on the ladder, not a failure of
            # ingestion — the next parser may handle this file.
            logger.info("extractor %s declined %s: %s", attempt.__name__, filename, e)
            continue
        if found is not None:
            return found

    raise DocumentError(
        f"Rule 5.3: could not extract text from {filename!r}: no parser "
        f"succeeded and the file is not valid UTF-8 text. It has not been "
        f"catalogued or indexed.")


def extract(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Bytes-in compatibility shape (the text routes and the tests): spool
    to a temp file and take the streamed path."""
    ext = os.path.splitext(filename)[1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext or ".bin") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        return extract_path(tmp_path, filename)
    finally:
        os.unlink(tmp_path)


def _docling(path: str, ext: str) -> Optional[ExtractionResult]:
    """Structured Markdown for PDF, DOCX, PPTX, XLSX, HTML and MD."""
    result = _converter().convert(path)
    text = result.document.export_to_markdown()
    return _result(text, "docling") if _ok(text) else None


def _pdf(path: str, ext: str) -> Optional[ExtractionResult]:
    if ext != ".pdf":
        return None
    for module_name, label in (("pypdf", "pypdf"), ("PyPDF2", "PyPDF2")):
        try:
            module = __import__(module_name)
            with open(path, "rb") as fh:
                reader = module.PdfReader(fh)
                pages = [p.extract_text() for p in reader.pages if p.extract_text()]
            text = "\n\n".join(pages)
            if _ok(text):
                return _result(text, label)
        except Exception as e:
            logger.info("%s extraction note: %s", label, e)
    return None


def _docx(path: str, ext: str) -> Optional[ExtractionResult]:
    if ext not in (".docx", ".doc"):
        return None
    import docx

    document = docx.Document(path)
    text = "\n\n".join(p.text for p in document.paragraphs if p.text)
    return _result(text, "python-docx") if _ok(text) else None


def _pptx(path: str, ext: str) -> Optional[ExtractionResult]:
    if ext not in (".pptx", ".ppt"):
        return None
    import pptx

    presentation = pptx.Presentation(path)
    slides = []
    for index, slide in enumerate(presentation.slides, 1):
        runs = [shape.text.strip() for shape in slide.shapes
                if hasattr(shape, "text") and shape.text]
        if runs:
            # The slide marker is kept: a chunk that lands mid-deck is far more
            # useful when it can say which slide it came from.
            slides.append(f"--- Slide {index} ---\n" + "\n".join(runs))
    text = "\n\n".join(slides)
    return _result(text, "python-pptx") if _ok(text) else None


def _xlsx(path: str, ext: str) -> Optional[ExtractionResult]:
    if ext not in (".xlsx", ".xls"):
        return None
    import openpyxl

    # read_only streams rows instead of materialising every sheet
    workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
    try:
        sheets = []
        for name in workbook.sheetnames:
            rows = []
            for row in workbook[name].iter_rows(values_only=True):
                line = " | ".join(str(v) for v in row if v is not None)
                if line.strip():
                    rows.append(line)
            if rows:
                sheets.append(f"--- Sheet: {name} ---\n" + "\n".join(rows))
    finally:
        workbook.close()
    text = "\n\n".join(sheets)
    return _result(text, "openpyxl") if _ok(text) else None


def _plain_text(path: str, ext: str) -> Optional[ExtractionResult]:
    """Strict UTF-8 (Rule 5.2), decoded incrementally in 1MB chunks.

    Not `errors="ignore"`: ignoring decode errors turns any binary file into
    mojibake that passes a non-empty check, so a PDF whose parsers had all
    failed was indexed as garbage rather than rejected.
    """
    decoder = codecs.getincrementaldecoder("utf-8")()
    parts = []
    try:
        with open(path, "rb") as fh:
            for chunk in iter(lambda: fh.read(1 << 20), b""):
                parts.append(decoder.decode(chunk))
            parts.append(decoder.decode(b"", final=True))
    except UnicodeDecodeError:
        return None
    text = "".join(parts)
    return _result(text, "utf8_text_reader") if _ok(text) else None
